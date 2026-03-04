from pptx import Presentation
import os

def extract_text_from_pptx(path):
    if not os.path.exists(path):
        print(f"File {path} not found.")
        return
    
    prs = Presentation(path)
    content = []
    
    for i, slide in enumerate(prs.slides):
        content.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text.strip())
        content.append("\n")
    
    return "\n".join(content)

if __name__ == "__main__":
    ppt_path = "updated Matrix - credentials. - Copy.pptx"
    text_content = extract_text_from_pptx(ppt_path)
    if text_content:
        with open("ppt_content.txt", "w", encoding="utf-8") as f:
            f.write(text_content)
        print("Successfully extracted content to ppt_content.txt")
